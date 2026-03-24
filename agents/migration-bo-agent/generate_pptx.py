"""
Generate a comprehensive BO-to-Fabric Migration PowerPoint deck.
Run: python generate_pptx.py
Output: BO_to_Fabric_Migration.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette (Microsoft Fabric inspired) ──────────────────────────
BG_DARK    = RGBColor(0x1B, 0x1A, 0x2D)   # dark navy
BG_MED     = RGBColor(0x24, 0x23, 0x3B)   # medium navy
ACCENT1    = RGBColor(0x00, 0x78, 0xD4)   # Microsoft blue
ACCENT2    = RGBColor(0x10, 0x7C, 0x10)   # green
ACCENT3    = RGBColor(0xF7, 0x63, 0x0C)   # orange
ACCENT_RED = RGBColor(0xD1, 0x34, 0x38)   # red
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY   = RGBColor(0x88, 0x88, 0x88)
GOLD       = RGBColor(0xFF, 0xB9, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper functions ────────────────────────────────────────────────────

def dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = spacing
        p.level = 0
    return txBox

def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Segoe UI"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Header row background
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ACCENT1
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_MED if r % 2 == 0 else BG_DARK
    return table_shape

def add_rounded_box(slide, left, top, width, height, fill_color, text="",
                    font_size=14, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(4)
    return shape

def section_divider(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    dark_bg(slide)
    add_text_box(slide, 1, 2.5, 11, 1.2, title, font_size=40, bold=True, color=ACCENT1,
                 alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, 2, 3.8, 9, 0.8, subtitle, font_size=20, color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.5),
                                   Inches(2.3), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT1
    line.line.fill.background()
    return slide


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_text_box(slide, 1.5, 1.8, 10, 1.5,
             "SAP BusinessObjects → Microsoft Fabric",
             font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 2, 3.3, 9, 0.8,
             "Migration Framework & Best Practices",
             font_size=26, color=ACCENT1, alignment=PP_ALIGN.CENTER)

# Accent line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.2),
                               Inches(3.3), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT1
line.line.fill.background()

add_text_box(slide, 2, 4.8, 9, 1,
             "Architecture Design  •  Migration Playbook  •  Risk Mitigation",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 2, 6.2, 9, 0.5,
             "Confidential — For Customer Discussion",
             font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_text_box(slide, 0.8, 0.4, 6, 0.8, "Agenda", font_size=32, bold=True, color=WHITE)

phases = [
    ("01", "Why Migrate?",             "Business drivers & SAP BO end-of-life reality", ACCENT1),
    ("02", "Discovery & Inventory",    "Asset mapping, usage analysis, complexity scoring", ACCENT1),
    ("03", "Architecture Patterns",    "4 migration patterns — from lift-and-shift to modernization", ACCENT1),
    ("04", "Component Mapping",        "Universe → Semantic Model, Webi → PBI, Crystal → Paginated", ACCENT1),
    ("05", "Security & Governance",    "CMC → Entra ID, row-level security, workspace topology", ACCENT1),
    ("06", "Common Pitfalls",          "15+ documented traps and how to avoid them", ACCENT_RED),
    ("07", "Migration Execution",      "Wave planning, parallel run, validation", ACCENT1),
    ("08", "Timeline & Estimation",    "Sizing model by perimeter complexity", ACCENT1),
    ("09", "Decommission & Adoption",  "BO shutdown checklist, user onboarding", ACCENT2),
]
for i, (num, title, desc, color) in enumerate(phases):
    y = 1.4 + i * 0.63
    add_rounded_box(slide, 0.8, y, 0.7, 0.5, color, num, font_size=16)
    add_text_box(slide, 1.65, y - 0.05, 4, 0.35, title, font_size=16, bold=True, color=WHITE)
    add_text_box(slide, 1.65, y + 0.25, 5.5, 0.35, desc, font_size=11, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHY MIGRATE
# ════════════════════════════════════════════════════════════════════════
section_divider("01 — Why Migrate?", "Business drivers & strategic context")

slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "Why Migrate from SAP BusinessObjects?", font_size=28, bold=True)

# Left column — Pain points
add_text_box(slide, 0.8, 1.3, 5.5, 0.5, "Current BO Challenges", font_size=18, bold=True, color=ACCENT_RED)
add_bullet_list(slide, 0.8, 1.9, 5.5, 4, [
    "⛔  SAP BI Platform approaching end of mainstream maintenance",
    "⛔  Flash-based components (Xcelsius/Design Studio) already EOL",
    "⛔  High licensing costs with limited innovation",
    "⛔  No native cloud, no AI/ML capabilities",
    "⛔  Complex administration (CMS, FRS, SIA, web servers)",
    "⛔  Poor mobile experience",
    "⛔  Self-service BI severely limited",
    "⛔  Talent pool shrinking (BO expertise declining)",
], font_size=13, color=LIGHT_GRAY)

# Right column — Fabric benefits
add_text_box(slide, 7, 1.3, 5.5, 0.5, "Microsoft Fabric Benefits", font_size=18, bold=True, color=ACCENT2)
add_bullet_list(slide, 7, 1.9, 5.5, 4, [
    "✅  Unified analytics platform (data + BI + AI)",
    "✅  OneLake — single data lake for all workloads",
    "✅  Direct Lake — near real-time without import overhead",
    "✅  Copilot AI integrated in reports & notebooks",
    "✅  Native cloud — zero infrastructure to manage",
    "✅  Self-service BI + governed enterprise BI",
    "✅  Excel, Teams, PowerPoint integration",
    "✅  Growing ecosystem & strong Microsoft investment",
], font_size=13, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 4-5 — DISCOVERY & INVENTORY
# ════════════════════════════════════════════════════════════════════════
section_divider("02 — Discovery & Inventory", "What do you have? What do you actually use?")

slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "BO Asset Inventory", font_size=28, bold=True)

data = [
    ["BO Asset", "Count", "Active?", "Fabric Target"],
    ["Universes (.unx / .unv)", "?", "?", "Semantic Model (Direct Lake / Import)"],
    ["Webi Reports (.wid)", "?", "?", "Power BI Reports (.pbix)"],
    ["Crystal Reports (.rpt)", "?", "?", "Paginated Reports (.rdl)"],
    ["Publications (bursting)", "?", "?", "PBI Subscriptions / Data-Driven"],
    ["Scheduled Instances", "?", "?", "Fabric Pipelines / Refresh"],
    ["Xcelsius / Design Studio", "?", "?", "Power BI Dashboards"],
    ["Live Office Documents", "?", "?", "Excel + Semantic Model"],
    ["BO SDK Custom Apps", "?", "?", "Fabric REST API / PBI Embedded"],
    ["CMC Users / Groups", "?", "?", "Entra ID + Workspace Roles + RLS"],
]
add_table(slide, 0.8, 1.3, 11.7, 5, len(data), 4, data, col_widths=[3.5, 1.2, 1.2, 5.8])

add_text_box(slide, 0.8, 6.5, 11, 0.5,
             "💡 Critical: Use BO Auditing database to identify active vs dormant reports. Retire unused reports — don't migrate them.",
             font_size=12, color=GOLD, bold=True)


# ── Complexity scoring slide ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "Complexity Scoring (per report/universe)", font_size=28, bold=True)

data = [
    ["Factor", "1 — Simple", "3 — Medium", "5 — Complex"],
    ["Data Sources", "Single table", "Multiple joins", "Cross-DB, derived tables, stored procs"],
    ["Prompts / Filters", "None", "Static LOV", "Cascading, optional prompts"],
    ["Formulas", "None", "Sum, Where", "RunningSum, Previous, Context ops"],
    ["Layout", "Simple table", "Charts + tables", "Cross-tabs, sections, breaks"],
    ["Security", "None", "Row-level in universe", "Dynamic row + profile-based"],
    ["Subreports", "None", "—", "Linked / embedded subreports"],
]
add_table(slide, 0.8, 1.3, 11.7, 4.5, len(data), 4, data, col_widths=[2.5, 3, 3, 3.2])

add_text_box(slide, 0.8, 6.2, 11, 0.6,
             "Score each report 1-5 per factor → Total complexity drives wave assignment and effort estimation",
             font_size=14, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 6-7 — ARCHITECTURE PATTERNS
# ════════════════════════════════════════════════════════════════════════
section_divider("03 — Architecture Patterns", "4 migration strategies — pick the right one for your context")

slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Migration Architecture Patterns", font_size=28, bold=True)

patterns = [
    ("A", "Lift & Shift", "Existing DB → Gateway → Semantic Model (Import)",
     "Fastest. Minimal data changes.", "Gateway dependency. No Fabric data benefits.", ACCENT3),
    ("B", "Lakehouse Modernization", "Sources → Pipelines → Lakehouse (Bronze/Silver/Gold) → Direct Lake",
     "Best performance. Future-proof.", "Requires data layer rebuild.", ACCENT2),
    ("C", "Warehouse Migration", "Sources → Pipelines → Fabric Warehouse (SQL) → Direct Lake",
     "Familiar SQL. Stored procs.", "Warehouse compute costs.", ACCENT1),
    ("D", "Hybrid / Phased", "Phase 1: Gateway → Phase 2: Lakehouse → Phase 3: Direct Lake",
     "Lowest risk. Incremental.", "Longer total timeline.", GOLD),
]

for i, (letter, name, flow, pros, cons, color) in enumerate(patterns):
    y = 1.4 + i * 1.45
    add_rounded_box(slide, 0.8, y, 0.6, 1.2, color, letter, font_size=20)
    add_text_box(slide, 1.55, y - 0.05, 4, 0.4, name, font_size=17, bold=True, color=WHITE)
    add_text_box(slide, 1.55, y + 0.35, 6, 0.35, flow, font_size=11, color=LIGHT_GRAY)
    add_text_box(slide, 1.55, y + 0.7, 5, 0.25, f"✅ {pros}", font_size=11, color=ACCENT2)
    add_text_box(slide, 1.55, y + 0.95, 5, 0.25, f"⚠️ {cons}", font_size=11, color=ACCENT3)

# Recommendation box
add_rounded_box(slide, 8, 1.5, 4.5, 5, BG_MED)
add_text_box(slide, 8.2, 1.6, 4.2, 0.5, "📌 Recommendation", font_size=16, bold=True, color=GOLD)
add_bullet_list(slide, 8.2, 2.2, 4.1, 4, [
    "Pattern B (Lakehouse) is recommended for most enterprise migrations",
    "Pattern D (Hybrid) is safest for risk-averse organizations",
    "Pattern A only if BO license expires imminently",
    "Pattern C for SQL-heavy teams with DW expertise",
], font_size=12, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 8-9 — COMPONENT MAPPING
# ════════════════════════════════════════════════════════════════════════
section_divider("04 — Component Mapping", "Every BO layer has a Fabric equivalent")

slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "BO → Fabric Component Mapping", font_size=28, bold=True)

data = [
    ["SAP BusinessObjects", "→", "Microsoft Fabric"],
    ["Universe (.unx / .unv)", "→", "Semantic Model (Direct Lake)"],
    ["  Business Layer", "→", "  DAX Measures & Calculated Columns"],
    ["  Data Foundation", "→", "  Lakehouse / Warehouse (SQL views)"],
    ["  Derived Tables", "→", "  Dataflows Gen2 / Lakehouse Views"],
    ["  @Prompt functions", "→", "  Slicers / Report Parameters"],
    ["  Row Restrictions", "→", "  Row-Level Security (RLS in DAX)"],
    ["  Aggregate Awareness", "→", "  Composite Models / Aggregations"],
    ["Webi Report (.wid)", "→", "Power BI Report (.pbix)"],
    ["Crystal Report (.rpt)", "→", "Paginated Report (.rdl)"],
    ["Publications (bursting)", "→", "Data-Driven Subscriptions"],
    ["CMC Scheduling", "→", "Fabric Pipelines / Refresh Schedules"],
    ["BI Launchpad / InfoView", "→", "Power BI Apps / Fabric Portal"],
    ["CMC Security", "→", "Entra ID + Workspace Roles + RLS"],
]
add_table(slide, 0.8, 1.3, 11.7, 5.8, len(data), 3, data, col_widths=[4.5, 0.7, 6.5])


# ── Formula translation ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "BO Formula → DAX Translation", font_size=28, bold=True)

data = [
    ["BO Formula", "DAX Equivalent"],
    ["=Sum([Revenue])", "Revenue = SUM(Sales[Revenue])"],
    ["=Sum([Revenue]) In [Region]", "CALCULATE(SUM(...), ALLEXCEPT(Sales, Geo[Region]))"],
    ["=RunningSum([Revenue])", "CALCULATE(SUM(...), FILTER(ALL(Dates), Dates[Date] <= MAX(...)))"],
    ["=Previous([Revenue])", "DATEADD / PARALLELPERIOD (time intelligence)"],
    ["=Percentage([Revenue])", "DIVIDE(SUM(...), CALCULATE(SUM(...), ALL(Sales)))"],
    ["=Where conditions", "CALCULATE with filter arguments"],
    ["NoFilter()", "ALL() / REMOVEFILTERS()"],
    ["ForEach / ForAll context", "CALCULATE + ALL / ALLEXCEPT (⚠️ complex rewrite)"],
]
add_table(slide, 0.8, 1.3, 11.7, 4.8, len(data), 2, data, col_widths=[5, 6.7])

add_text_box(slide, 0.8, 6.3, 11, 0.6,
             "⚠️ BO Context Operators (In, ForEach, ForAll) have NO direct DAX equivalent — they require expert DAX rewriting and thorough testing",
             font_size=13, color=ACCENT_RED, bold=True)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 10-11 — SECURITY
# ════════════════════════════════════════════════════════════════════════
section_divider("05 — Security & Governance", "CMC → Entra ID + Fabric RBAC + RLS")

slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "Security Model Migration", font_size=28, bold=True)

data = [
    ["BO CMC Security", "Fabric Equivalent"],
    ["Enterprise Users (AD sync)", "Entra ID Users (native sync)"],
    ["BO Groups", "Entra ID Security Groups"],
    ["Folder: View rights", "Workspace Viewer role"],
    ["Folder: Schedule rights", "Workspace Contributor role"],
    ["Folder: Edit rights", "Workspace Member role"],
    ["Folder: Full Control", "Workspace Admin role"],
    ["Universe Row Restriction", "RLS Role + DAX filter"],
    ["Object-level security (hide columns)", "OLS (Object-Level Security) in Semantic Model"],
    ["Connection-level access", "Gateway data source permissions"],
    ["BO Profiles (dynamic security)", "Dynamic RLS with USERPRINCIPALNAME()"],
]
add_table(slide, 0.8, 1.3, 11.7, 5.5, len(data), 2, data, col_widths=[5, 6.7])

# Workspace topology
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "Recommended Workspace Topology", font_size=28, bold=True)

workspaces = [
    ("📊", "[CUSTOMER]-DataPlatform", "Lakehouse, Warehouse, Pipelines", ACCENT1),
    ("🧮", "[CUSTOMER]-SemanticModels", "Shared models (replaces Universes)", ACCENT1),
    ("💰", "[CUSTOMER]-Finance-Reports", "Finance domain reports", ACCENT2),
    ("📈", "[CUSTOMER]-Sales-Reports", "Sales domain reports", ACCENT2),
    ("👥", "[CUSTOMER]-HR-Reports", "HR reports (⚠️ RLS critical)", ACCENT_RED),
    ("📋", "[CUSTOMER]-PaginatedReports", "Migrated Crystal Reports", ACCENT3),
    ("🔧", "[CUSTOMER]-Dev", "Development / staging", MED_GRAY),
]

for i, (icon, name, desc, color) in enumerate(workspaces):
    y = 1.4 + i * 0.75
    add_rounded_box(slide, 1, y, 0.6, 0.55, color, icon, font_size=18)
    add_text_box(slide, 1.8, y - 0.02, 5, 0.35, name, font_size=15, bold=True, color=WHITE)
    add_text_box(slide, 1.8, y + 0.3, 5, 0.3, desc, font_size=12, color=LIGHT_GRAY)

# Governance box
add_rounded_box(slide, 7.5, 1.4, 5, 5, BG_MED)
add_text_box(slide, 7.7, 1.5, 4.7, 0.5, "Governance Principles", font_size=16, bold=True, color=GOLD)
add_bullet_list(slide, 7.7, 2.1, 4.5, 4, [
    "Separate data, models, and reports into different workspaces",
    "Shared Semantic Models workspace = new 'Universe' layer",
    "Each domain owns its report workspace",
    "Dev workspace for testing before promotion",
    "Entra ID groups control access — no individual permissions",
    "RLS mandatory for sensitive data (HR, Finance)",
    "Deployment pipelines for CI/CD (Dev → Test → Prod)",
], font_size=12, color=LIGHT_GRAY, spacing=Pt(8))


# ════════════════════════════════════════════════════════════════════════
# SLIDE 12-13 — COMMON PITFALLS
# ════════════════════════════════════════════════════════════════════════
section_divider("06 — Common Pitfalls", "15+ traps to avoid in BO → Fabric migrations")

# Architecture pitfalls
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "Architecture & Organizational Pitfalls", font_size=28, bold=True, color=ACCENT_RED)

data = [
    ["#", "Pitfall", "Impact", "Mitigation"],
    ["1", "1:1 report migration", "Wasted effort on unused reports", "Use BO Auditing — retire dormant first"],
    ["2", "No shared Semantic Model", "Report sprawl, inconsistent metrics", "Build shared models = new Universes"],
    ["3", "Gateway bottleneck", "Refresh failures, performance", "Plan gateway capacity or go Lakehouse"],
    ["4", "Underestimate cross-tabs", "Matrix visuals ≠ Webi cross-tabs", "Prototype complex cross-tabs early"],
    ["5", "No parallel run", "User distrust of new reports", "Run BO + Fabric side-by-side 2-4 weeks"],
    ["6", "Big bang migration", "Risk of failure, user rejection", "Wave-based: department by department"],
    ["7", "No user training", "Support flood: 'Where's my report?'", "Train power users first, create videos"],
    ["8", "Losing BO audit trail", "Compliance gap", "Export BO audit data before decommission"],
]
add_table(slide, 0.5, 1.3, 12.3, 5.5, len(data), 4, data, col_widths=[0.5, 2.8, 3.5, 5.5])

# Technical pitfalls
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "Technical Pitfalls", font_size=28, bold=True, color=ACCENT_RED)

data = [
    ["#", "Pitfall", "Impact", "Mitigation"],
    ["9", "Context operators (In, ForEach, ForAll)", "No direct DAX equivalent", "Deep DAX rewrite + test with sample data"],
    ["10", "@Prompt cascading", "Slicers don't cascade same way", "Field parameters or manage via relationships"],
    ["11", "Derived tables with complex SQL", "Can't just copy SQL", "Rewrite as Lakehouse views / Dataflow M"],
    ["12", "Crystal embedded formulas", "Some have no RDL equivalent", "Manual rewrite; keep Crystal if low priority"],
    ["13", "RLS security gaps", "BO dynamic profile ≠ Fabric RLS", "Full security audit before go-live"],
    ["14", "Date/time handling", "BO date functions ≠ DAX", "Build proper date dimension from day 1"],
    ["15", "Large dataset limits", "PBI Import has visual limits", "Use Direct Lake for >100M rows"],
    ["16", "NULL vs BLANK()", "Rounding & NULL differences", "Test edge cases, use ROUND() explicitly"],
]
add_table(slide, 0.5, 1.3, 12.3, 5.5, len(data), 4, data, col_widths=[0.5, 3.5, 3, 5.3])


# ════════════════════════════════════════════════════════════════════════
# SLIDE 14-15 — MIGRATION EXECUTION
# ════════════════════════════════════════════════════════════════════════
section_divider("07 — Migration Execution", "Wave planning, validation, parallel run")

slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "Wave-Based Migration Approach", font_size=28, bold=True)

waves = [
    ("Wave 0", "Foundation", "Lakehouse/Warehouse setup, Semantic Models (Universes), Security, Gateway",
     "2-4 weeks", ACCENT1),
    ("Wave 1", "Quick Wins", "Top 10 most-used simple reports (complexity 1-2)",
     "2-4 weeks", ACCENT2),
    ("Wave 2", "Department A", "First full department — all reports, Crystal, publications",
     "4-8 weeks", ACCENT1),
    ("Wave 3", "Department B", "Second department — apply lessons learned from Wave 2",
     "4-6 weeks", ACCENT1),
    ("Wave N", "Remaining", "All other departments + complex reports + SDK apps",
     "Variable", MED_GRAY),
    ("Final", "Decommission", "Parallel run complete → BO shutdown",
     "2-4 weeks", ACCENT_RED),
]

for i, (wave, name, desc, duration, color) in enumerate(waves):
    y = 1.4 + i * 0.95
    add_rounded_box(slide, 0.8, y, 1.6, 0.7, color, wave, font_size=13)
    add_text_box(slide, 2.6, y - 0.05, 3.5, 0.35, name, font_size=16, bold=True, color=WHITE)
    add_text_box(slide, 2.6, y + 0.3, 6, 0.35, desc, font_size=11, color=LIGHT_GRAY)
    add_text_box(slide, 10.5, y + 0.1, 2, 0.35, duration, font_size=13, bold=True, color=GOLD,
                 alignment=PP_ALIGN.RIGHT)


# Validation checklist
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "Validation Checklist (per report)", font_size=28, bold=True)

add_bullet_list(slide, 0.8, 1.4, 5.5, 5, [
    "☐  Data accuracy: BO output vs Fabric output (same date, filters)",
    "☐  Row counts match between BO and Fabric",
    "☐  Aggregations match (SUM, AVG, COUNT, DISTINCT)",
    "☐  Filters/slicers = same results as BO prompts",
    "☐  RLS: Same user sees same data scope",
    "☐  Schedule: Refresh runs at expected times",
    "☐  Performance: Page load < 5 seconds",
    "☐  Export: PDF/Excel output matches BO format",
    "☐  Mobile: Renders correctly (if required)",
], font_size=14, color=LIGHT_GRAY)

# Parallel run protocol
add_text_box(slide, 7, 1.3, 5.5, 0.5, "Parallel Run Protocol", font_size=18, bold=True, color=GOLD)
add_bullet_list(slide, 7, 1.9, 5.5, 4, [
    "1.  Run BO + Fabric side-by-side for 2-4 weeks",
    "2.  Business users compare 5-10 critical reports daily",
    "3.  Log discrepancies in shared tracker",
    "4.  Fix and re-validate each discrepancy",
    "5.  Business sign-off per wave before BO cutover",
    "",
    "⚠️  Never skip parallel run — it is the #1 confidence builder for end users",
], font_size=13, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 16 — TIMELINE & ESTIMATION
# ════════════════════════════════════════════════════════════════════════
section_divider("08 — Timeline & Estimation", "Right-size the effort to the perimeter")

slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "Estimation Model", font_size=28, bold=True)

data = [
    ["Activity", "Small (<50 reports)", "Medium (50-200)", "Large (200+)"],
    ["Discovery & Inventory", "1-2 weeks", "2-3 weeks", "3-4 weeks"],
    ["Architecture Design", "1 week", "2 weeks", "2-3 weeks"],
    ["Semantic Model Build", "2-4 weeks", "4-8 weeks", "8-16 weeks"],
    ["Webi → PBI Report Migration", "4-8 weeks", "8-20 weeks", "20-40+ weeks"],
    ["Crystal → Paginated Reports", "2-4 weeks", "4-10 weeks", "10-20 weeks"],
    ["Security Setup & Testing", "1 week", "2-3 weeks", "3-4 weeks"],
    ["Testing & Parallel Run", "2-3 weeks", "3-4 weeks", "4-6 weeks"],
    ["Training & Adoption", "1 week", "2 weeks", "3-4 weeks"],
    ["Total Estimated Duration", "3-5 months", "6-12 months", "12-24 months"],
]
add_table(slide, 0.8, 1.3, 11.7, 5.5, len(data), 4, data, col_widths=[4, 2.6, 2.6, 2.5])

add_text_box(slide, 0.8, 6.8, 11, 0.4,
             "Rule of thumb:  ~0.5-2 days/Webi report  •  ~1-3 days/Crystal Report  •  ~1-2 weeks/Universe → Semantic Model",
             font_size=12, color=GOLD, bold=True)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 17 — DECOMMISSION & ADOPTION
# ════════════════════════════════════════════════════════════════════════
section_divider("09 — Decommission & Adoption", "Shut down BO safely, onboard users confidently")

slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "Decommission & Adoption", font_size=28, bold=True)

# Left — decommission
add_text_box(slide, 0.8, 1.3, 5.5, 0.5, "BO Decommission Checklist", font_size=18, bold=True, color=ACCENT_RED)
add_bullet_list(slide, 0.8, 1.9, 5.5, 4.5, [
    "☐  All active reports migrated and validated",
    "☐  BO audit data exported and archived",
    "☐  Universe documentation archived",
    "☐  BO connection details documented",
    "☐  BO server backup taken before shutdown",
    "☐  DNS / bookmarks redirected to Fabric",
    "☐  BO license termination communicated to SAP",
    "☐  BO servers decommissioned from infrastructure",
], font_size=13, color=LIGHT_GRAY)

# Right — adoption
add_text_box(slide, 7, 1.3, 5.5, 0.5, "Adoption Accelerators", font_size=18, bold=True, color=ACCENT2)
add_bullet_list(slide, 7, 1.9, 5.5, 4.5, [
    "✅  Power BI App per department (replaces BI Launchpad)",
    "✅  Quick-reference card: 'BO → Fabric: Where's my report?'",
    "✅  Power user champions network",
    "✅  Video walkthroughs for top 20 reports",
    "✅  Office hours — first 2 weeks post go-live",
    "✅  Feedback loop for continuous improvement",
    "✅  Copilot demos to showcase AI capabilities",
    "✅  Monthly adoption dashboard (usage metrics)",
], font_size=13, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 18 — KEY TAKEAWAYS
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Key Takeaways", font_size=32, bold=True, color=GOLD)

takeaways = [
    ("1", "Retire Before Migrating", "Use BO audit data to kill dormant reports. Migrate 60-70%, retire the rest.", ACCENT1),
    ("2", "Semantic Layer First", "Build shared Semantic Models before reports — they replace Universes.", ACCENT1),
    ("3", "Wave-Based Delivery", "Department by department. Never big bang.", ACCENT1),
    ("4", "Parallel Run is Mandatory", "2-4 weeks BO + Fabric side-by-side. Business sign-off before cutover.", ACCENT_RED),
    ("5", "Invest in DAX Expertise", "Context operators & complex BO formulas require skilled DAX developers.", ACCENT3),
    ("6", "Security Audit Before Go-Live", "Full RLS testing. Same user must see same data in both platforms.", ACCENT_RED),
]

for i, (num, title, desc, color) in enumerate(takeaways):
    y = 1.5 + i * 0.95
    add_rounded_box(slide, 0.8, y, 0.6, 0.7, color, num, font_size=18)
    add_text_box(slide, 1.6, y - 0.02, 5, 0.35, title, font_size=17, bold=True, color=WHITE)
    add_text_box(slide, 1.6, y + 0.35, 10, 0.35, desc, font_size=13, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 19 — CLOSING
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_text_box(slide, 1.5, 2.2, 10, 1.2,
             "Next Steps",
             font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.4),
                               Inches(2.3), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT1
line.line.fill.background()

add_bullet_list(slide, 3, 3.8, 7, 3, [
    "1.  Complete the Discovery Questionnaire with BO Admin team",
    "2.  Extract BO Auditing data for usage-based prioritization",
    "3.  Score top 50 reports for complexity → define Wave 1 scope",
    "4.  Architecture decision: Pattern A / B / C / D",
    "5.  Build POC: 3-5 reports end-to-end (Universe → Semantic Model → Report)",
], font_size=16, color=LIGHT_GRAY, spacing=Pt(12))

add_text_box(slide, 2, 6.5, 9, 0.5,
             "Thank You",
             font_size=20, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════════
import os
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "BO_to_Fabric_Migration.pptx")
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
